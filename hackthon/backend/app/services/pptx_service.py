from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def generate_pptx(
    ai_result: dict[str, Any],
    output_path: Path,
) -> Path:
    """
    Convert Qwen-generated slide JSON into a PowerPoint presentation.

    Expected structure:

    {
        "title": "...",
        "summary": "...",
        "slides": [
            {
                "slide_number": 1,
                "layout": "title",
                "title": "...",
                "bullets": [...]
            }
        ]
    }
    """

    if not ai_result:
        raise ValueError("AI result is empty.")

    slides = ai_result.get("slides", [])

    if not slides:
        raise ValueError("AI result contains no slides.")

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()

    # Use a widescreen presentation.
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    # Remove the default slide created by Presentation().
    if len(presentation.slides) > 0:
        first_slide = presentation.slides[0]
        r_id = first_slide.slide_id
        presentation.part.drop_rel(r_id)
        presentation.slides._sldIdLst.remove(
            presentation.slides._sldIdLst[0]
        )

    for slide_data in slides:
        layout = str(slide_data.get("layout", "overview")).lower()
        title = str(slide_data.get("title", "Untitled Slide"))
        bullets = slide_data.get("bullets", [])

        if layout == "title":
            slide_layout = presentation.slide_layouts[0]
        else:
            slide_layout = presentation.slide_layouts[1]

        slide = presentation.slides.add_slide(slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(8, 15, 30)

        # Title
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title

            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(30)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(
                    255,
                    255,
                    255,
                )

        # Title slide
        if layout == "title":
            subtitle_text = ai_result.get("summary", "")

            if subtitle_text:
                subtitle = slide.shapes.add_textbox(
                    Inches(1.0),
                    Inches(3.2),
                    Inches(11.3),
                    Inches(1.5),
                )

                text_frame = subtitle.text_frame
                text_frame.word_wrap = True
                text_frame.text = str(subtitle_text)

                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(20)
                    paragraph.font.color.rgb = RGBColor(
                        180,
                        190,
                        205,
                    )

            continue

        # Content slides
        if slide.placeholders:
            body = None

            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:
                    body = shape
                    break

            if body is not None:
                text_frame = body.text_frame
                text_frame.clear()

                for index, bullet in enumerate(bullets):
                    paragraph = (
                        text_frame.paragraphs[0]
                        if index == 0
                        else text_frame.add_paragraph()
                    )

                    paragraph.text = str(bullet)
                    paragraph.level = 0
                    paragraph.font.size = Pt(20)
                    paragraph.font.color.rgb = RGBColor(
                        220,
                        225,
                        235,
                    )
                    paragraph.space_after = Pt(12)

    presentation.save(output_path)

    return output_path